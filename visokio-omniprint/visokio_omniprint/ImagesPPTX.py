import os
from .DriverBase import DriverBase
from pptx import Presentation
from pptx.util import Inches

class ImagesPPTX(DriverBase):


    def __init__(self):
        self.paths = []


    def add_path(self, path):
        self.paths.append(path)


    def _split_resolution(self, orientation, resolution):
        resSplit = resolution.split("x")
        #width, then height
        imgSize = [ int(resSplit[0]), int(resSplit[1]) ]

        if orientation == "P":
            resolution = resSplit[1]+"x"+resSplit[0]
            imgSize = [ int(resSplit[1]), int(resSplit[0]) ]

        #Pdf lib wants heigth then width
        resSplit[0], resSplit[1] = resSplit[1], resSplit[0]

        return resSplit

    def create_pptx_in_path(self, path, orientation, resolution):
        res = self._split_resolution(orientation, resolution)

        prs = Presentation()

        for image_path in self.paths:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slide.shapes.add_picture(image_path, left, top, Inches(10))

        if os.path.exists(path):
            os.remove(path)
        prs.save(path)
